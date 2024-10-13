import os
import csv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_ppt_content(ppt_file, output_csv, image_dir):
    # Create a directory to save extracted images
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    # Open the PowerPoint presentation
    prs = Presentation(ppt_file)

    # Prepare to write to CSV with utf-8 encoding
    with open(output_csv, mode='w', newline='', encoding='utf-8') as csvfile:
        fieldnames = ['Slide Number', 'Image Path', 'Text', 'Notes']
        writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
        writer.writeheader()

        # Loop through each slide
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            image_paths = []
            slide_notes = []

            # Extract text from the slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text.append(paragraph.text)

                # Extract images from the slide
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_name = f"slide_{slide_num}_image_{len(image_paths) + 1}.png"
                    image_path = os.path.join(image_dir, image_name)
                    with open(image_path, 'wb') as f:
                        f.write(image.blob)
                    image_paths.append(image_path)

            # Extract notes from the slide
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                if notes_text_frame:
                    for paragraph in notes_text_frame.paragraphs:
                        slide_notes.append(paragraph.text)

            # Write slide information to CSV, with None for empty fields
            writer.writerow({
                'Slide Number': slide_num,
                'Image Path': "; ".join(image_paths) if image_paths else 'None',
                'Text': " ".join(slide_text) if slide_text else 'None',
                'Notes': " ".join(slide_notes) if slide_notes else 'None'
            })

# Main script execution
if __name__ == "__main__":
    # Prompt user for the PowerPoint file path
    ppt_file = input("Please enter the path to the PowerPoint file (e.g., C:\\path\\to\\file.pptx): ")

    # Check if the provided path is valid
    if os.path.isfile(ppt_file):
        output_csv = 'output.csv'  # CSV output file path
        image_dir = 'images'  # Directory to save images

        # Call the function to extract content
        extract_ppt_content(ppt_file, output_csv, image_dir)
        print(f"Content extracted to '{output_csv}' and images saved in '{image_dir}'.")
    else:
        print("Invalid file path. Please make sure the file exists and try again.")
